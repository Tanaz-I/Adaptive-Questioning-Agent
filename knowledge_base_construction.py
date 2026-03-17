"""
PPTX RAG Pipeline
=================
Reads all .pptx files from a directory, extracts content slide-by-slide,
chunks and vectorizes the text, and stores it in a ChromaDB vector database
for use in a Retrieval-Augmented Generation (RAG) pipeline.

Dependencies:
    pip install markitdown[pptx] python-pptx chromadb sentence-transformers tqdm
"""

import os
import re
import json
import hashlib
from pathlib import Path
from dataclasses import dataclass, field, asdict
from typing import Optional

from pptx import Presentation
from markitdown import MarkItDown
import chromadb
from chromadb.config import Settings
from sentence_transformers import SentenceTransformer
from tqdm import tqdm


# ─────────────────────────────────────────────
# Configuration
# ─────────────────────────────────────────────

PPTX_DIR        = "./presentations"          # Directory containing .pptx files
CHROMA_DB_DIR   = "./chroma_db"              # Persistent ChromaDB storage path
COLLECTION_NAME = "pptx_rag"                 # ChromaDB collection name
EMBED_MODEL     = "all-MiniLM-L6-v2"        # Sentence-Transformers model
CHUNK_SIZE      = 500                        # Max characters per chunk
CHUNK_OVERLAP   = 80                         # Overlap between consecutive chunks


# ─────────────────────────────────────────────
# Data model
# ─────────────────────────────────────────────

@dataclass
class SlideChunk:
    """A single text chunk extracted from one slide."""
    doc_id:        str            # Unique ID for this chunk
    text:          str            # The chunk text
    file_name:     str            # Source .pptx filename
    file_path:     str            # Absolute path to source file
    slide_number:  int            # 1-based slide index
    slide_title:   str            # Title of the slide (if present)
    chunk_index:   int            # Index of this chunk within the slide
    total_chunks:  int            # Total chunks from this slide
    speaker_notes: str = ""       # Speaker notes from the slide
    extra_metadata: dict = field(default_factory=dict)


# ─────────────────────────────────────────────
# PPTX Extraction
# ─────────────────────────────────────────────

def _slide_title(slide) -> str:
    """Return the title placeholder text, or an empty string."""
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        return slide.shapes.title.text_frame.text.strip()
    return ""


def _slide_body_text(slide) -> str:
    """
    Concatenate all text from non-title placeholders and free text boxes.
    Uses python-pptx for structured extraction (preserves bullet hierarchy).
    """
    lines = []
    title_shape = slide.shapes.title
    for shape in slide.shapes:
        if shape == title_shape:
            continue
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    level = para.level  # indent level (0 = top)
                    indent = "  " * level
                    lines.append(f"{indent}{text}")
    return "\n".join(lines)


def _slide_notes(slide) -> str:
    """Extract speaker notes."""
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        return tf.text.strip() if tf else ""
    except Exception:
        return ""


def extract_slides(pptx_path: Path) -> list[dict]:
    """
    Parse a .pptx file and return a list of slide dicts:
        {slide_number, title, body, notes, markdown}
    """
    prs = Presentation(str(pptx_path))
    md_converter = MarkItDown()

    # Full-file markdown as a fallback / supplement
    full_md = md_converter.convert(str(pptx_path)).text_content

    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        title  = _slide_title(slide)
        body   = _slide_body_text(slide)
        notes  = _slide_notes(slide)

        slides.append({
            "slide_number": idx,
            "title":        title,
            "body":         body,
            "notes":        notes,
        })

    return slides


# ─────────────────────────────────────────────
# Chunking
# ─────────────────────────────────────────────

def _chunk_text(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """
    Split `text` into overlapping fixed-size character chunks.
    Tries to break on sentence boundaries first.
    """
    if not text or len(text) <= size:
        return [text] if text.strip() else []

    # Split on sentence endings
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks = []
    current = ""

    for sentence in sentences:
        if len(current) + len(sentence) + 1 <= size:
            current = (current + " " + sentence).strip()
        else:
            if current:
                chunks.append(current)
            # If a single sentence exceeds size, hard-split it
            while len(sentence) > size:
                chunks.append(sentence[:size])
                sentence = sentence[size - overlap:]
            current = sentence

    if current:
        chunks.append(current)

    return chunks


def slide_to_chunks(slide: dict, file_name: str, file_path: str) -> list[SlideChunk]:
    """Convert a single slide dict into one or more SlideChunk objects."""
    # Compose the full text for this slide
    parts = []
    if slide["title"]:
        parts.append(f"Title: {slide['title']}")
    if slide["body"]:
        parts.append(slide["body"])
    if slide["notes"]:
        parts.append(f"Speaker Notes: {slide['notes']}")

    full_text = "\n\n".join(parts).strip()
    raw_chunks = _chunk_text(full_text)

    result = []
    for chunk_idx, chunk_text in enumerate(raw_chunks):
        # Build a stable, deterministic ID
        unique_str = f"{file_path}::slide{slide['slide_number']}::chunk{chunk_idx}"
        doc_id = hashlib.md5(unique_str.encode()).hexdigest()

        result.append(SlideChunk(
            doc_id        = doc_id,
            text          = chunk_text,
            file_name     = file_name,
            file_path     = file_path,
            slide_number  = slide["slide_number"],
            slide_title   = slide["title"],
            chunk_index   = chunk_idx,
            total_chunks  = len(raw_chunks),
            speaker_notes = slide["notes"],
        ))

    return result


# ─────────────────────────────────────────────
# Vector DB (ChromaDB)
# ─────────────────────────────────────────────

def build_vector_db(
    chunks:          list[SlideChunk],
    embed_model:     str = EMBED_MODEL,
    chroma_dir:      str = CHROMA_DB_DIR,
    collection_name: str = COLLECTION_NAME,
) -> chromadb.Collection:
    """
    Embed all chunks and upsert them into a persistent ChromaDB collection.
    Returns the collection object.
    """
    print(f"\n[VectorDB] Loading embedding model: {embed_model}")
    model = SentenceTransformer(embed_model)

    print(f"[VectorDB] Embedding {len(chunks)} chunks …")
    texts = [c.text for c in chunks]
    embeddings = model.encode(texts, show_progress_bar=True, batch_size=64)

    # Initialise persistent ChromaDB
    client = chromadb.PersistentClient(
        path=chroma_dir,
        settings=Settings(anonymized_telemetry=False),
    )
    collection = client.get_or_create_collection(
        name=collection_name,
        metadata={"hnsw:space": "cosine"},
    )

    # Upsert in batches of 500 (ChromaDB hard limit)
    BATCH = 500
    print(f"[VectorDB] Upserting into collection '{collection_name}' …")
    for i in tqdm(range(0, len(chunks), BATCH), desc="Batches"):
        batch = chunks[i : i + BATCH]
        collection.upsert(
            ids        = [c.doc_id for c in batch],
            documents  = [c.text   for c in batch],
            embeddings = embeddings[i : i + BATCH].tolist(),
            metadatas  = [
                {
                    "file_name":    c.file_name,
                    "file_path":    c.file_path,
                    "slide_number": c.slide_number,
                    "slide_title":  c.slide_title,
                    "chunk_index":  c.chunk_index,
                    "total_chunks": c.total_chunks,
                    "speaker_notes": c.speaker_notes[:500],   # truncate for metadata limit
                }
                for c in batch
            ],
        )

    print(f"[VectorDB] ✓ {collection.count()} total documents in collection.")
    return collection


# ─────────────────────────────────────────────
# RAG Query helper
# ─────────────────────────────────────────────

def query_rag(
    query:           str,
    collection:      chromadb.Collection,
    embed_model:     str  = EMBED_MODEL,
    top_k:           int  = 5,
) -> list[dict]:
    """
    Embed a query and retrieve the top-k most relevant chunks.

    Returns a list of dicts with keys:
        text, score, file_name, slide_number, slide_title, chunk_index
    """
    model = SentenceTransformer(embed_model)
    query_embedding = model.encode([query])[0].tolist()

    results = collection.query(
        query_embeddings=[query_embedding],
        n_results=top_k,
        include=["documents", "metadatas", "distances"],
    )

    hits = []
    for doc, meta, dist in zip(
        results["documents"][0],
        results["metadatas"][0],
        results["distances"][0],
    ):
        hits.append({
            "text":         doc,
            "score":        round(1 - dist, 4),   # cosine similarity
            "file_name":    meta["file_name"],
            "slide_number": meta["slide_number"],
            "slide_title":  meta["slide_title"],
            "chunk_index":  meta["chunk_index"],
        })

    return hits


# ─────────────────────────────────────────────
# Main pipeline
# ─────────────────────────────────────────────

def run_pipeline(
    pptx_dir:        str = PPTX_DIR,
    chroma_dir:      str = CHROMA_DB_DIR,
    collection_name: str = COLLECTION_NAME,
    embed_model:     str = EMBED_MODEL,
) -> chromadb.Collection:
    """
    End-to-end pipeline:
        1. Discover all .pptx files in `pptx_dir`
        2. Extract & chunk slide content
        3. Embed and store in ChromaDB
        4. Return the live collection for querying
    """
    import subprocess

    for ppt in Path(pptx_dir).glob("**/*.ppt"):
        subprocess.run([
            "libreoffice", "--headless", "--convert-to", "pptx",
            "--outdir", str(ppt.parent), str(ppt)
        ])

    pptx_files = sorted(Path(pptx_dir).glob("**/*.pptx"))
    if not pptx_files:
        raise FileNotFoundError(f"No .pptx files found in: {pptx_dir}")

    print(f"\n{'='*55}")
    print(f"  PPTX RAG Pipeline")
    print(f"{'='*55}")
    print(f"  Source dir : {pptx_dir}")
    print(f"  Found      : {len(pptx_files)} file(s)")
    print(f"  ChromaDB   : {chroma_dir}")
    print(f"  Collection : {collection_name}")
    print(f"  Embed model: {embed_model}")
    print(f"{'='*55}\n")

    all_chunks: list[SlideChunk] = []

    for pptx_path in tqdm(pptx_files, desc="Parsing files"):
        print(f"\n  → {pptx_path.name}")
        try:
            slides = extract_slides(pptx_path)
            file_chunks = []
            for slide in slides:
                file_chunks.extend(
                    slide_to_chunks(slide, pptx_path.name, str(pptx_path.resolve()))
                )
            print(f"     {len(slides)} slides  |  {len(file_chunks)} chunks")
            all_chunks.extend(file_chunks)
        except Exception as exc:
            print(f"     [WARN] Skipping {pptx_path.name}: {exc}")

    if not all_chunks:
        raise ValueError("No chunks were extracted. Check your .pptx files.")

    print(f"\n[Pipeline] Total chunks across all files: {len(all_chunks)}")

    collection = build_vector_db(
        chunks          = all_chunks,
        embed_model     = embed_model,
        chroma_dir      = chroma_dir,
        collection_name = collection_name,
    )

    print("\n[Pipeline] ✓ Done. Vector DB is ready for queries.\n")
    return collection


# ─────────────────────────────────────────────
# Entry point / demo
# ─────────────────────────────────────────────

if __name__ == "__main__":
    # ── 1. Run the ingestion pipeline ───────────────────────────────────────
    collection = run_pipeline()

    # ── 2. Demo: run a sample query ─────────────────────────────────────────
    sample_query = "What are the key conclusions from the presentation?"
    print(f"[Query] '{sample_query}'\n")

    hits = query_rag(sample_query, collection, top_k=3)
    for rank, hit in enumerate(hits, 1):
        print(f"  Rank {rank}  |  score={hit['score']}  |  "
              f"{hit['file_name']}  Slide {hit['slide_number']}  "
              f"(chunk {hit['chunk_index']})")
        print(f"  Title : {hit['slide_title'] or '—'}")
        print(f"  Text  : {hit['text'][:200]}{'…' if len(hit['text']) > 200 else ''}")
        print()