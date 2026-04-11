"""
Generalized Knowledge Base Construction

Reads all supported files from a directory, extracts text,
chunks it, vectorizes, and stores in ChromaDB for RAG.

Supported formats:
    .pptx / .ppt   — Presentation slides   (python-pptx)
    .pdf           — PDF documents         (pdfplumber)
    .docx / .doc   — Word documents        (python-docx)
    .txt           — Plain text files

Dependencies:
    pip install python-pptx pdfplumber python-docx \
                chromadb sentence-transformers tqdm
    pip install easyocr opencv-python pillow
    sudo apt install libreoffice -y   # for .ppt and .doc conversion

Usage:
    python knowledge_base_construction.py
"""

import re
import io
import base64
import shutil
import hashlib
import subprocess
from pathlib import Path
from dataclasses import dataclass
import json
import base64
import requests
import pdfplumber
from pptx import Presentation
from docx import Document as DocxDocument

from PIL import Image
import chromadb
from chromadb.config import Settings
from sentence_transformers import SentenceTransformer
from tqdm import tqdm
from sklearn.metrics.pairwise import cosine_similarity

from img2table.document import Image as Img2TableImage
from img2table.ocr import TesseractOCR
from transformers import Pix2StructProcessor, Pix2StructForConditionalGeneration
processor = Pix2StructProcessor.from_pretrained("google/deplot")
model = Pix2StructForConditionalGeneration.from_pretrained("google/deplot")

import nltk
nltk.download("punkt")
nltk.download("punkt_tab")
from nltk.tokenize import sent_tokenize
from NLP.image_processing import *

pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Configuration

DOCS_DIR         = "./contents"       # Root directory with all source files
CHROMA_DB_DIR    = "./chroma_db"
COLLECTION_NAME  = "rag_kb"
EMBED_MODEL      = "all-MiniLM-L6-v2"
MULTIMODAL_MODEL = "llava"

CHUNK_SIZE              = 500                 # Max characters per chunk
CHUNK_OVERLAP           = 80                  # Overlap between chunks
CHUNKING_STRATEGY       = "sentence"          # semantic or sentence

SEMANTIC_SIM_THRESHOLD  = 0.4

# Supported extensions
SUPPORTED_EXTS  = {".pptx", ".ppt", ".docx", ".doc", ".pdf", ".txt"}
# reader = easyocr.Reader(['en'])




@dataclass
class Chunk:
    doc_id:       str           # Unique MD5 ID
    text:         str           # Chunk text
    file_name:    str           # Source filename
    file_path:    str           # Absolute path
    file_type:    str           # pptx / pdf / docx / txt
    page_number:  int           # Page/slide number (1-based); 0 if N/A
    section:      str           # Section/heading/slide title if available
    chunk_index:  int           # Index within the page/section
    total_chunks: int           # Total chunks from this page/section
    is_image: bool
    image_type:   str           # image, table, chart, or text
    parent_id:    str = ""          # NEW — MD5 of (file_path + page_number)
    chunk_type:   str = "text"      # NEW — "text" | "code" | "equation" | "image"


# Legacy format conversion (.ppt / .doc -> modern)

def convert_legacy_files(docs_dir: str):
    """Convert old .ppt/.doc to .pptx/.docx using LibreOffice."""
    converted = 0
    for ext, target_ext in [(".ppt", "pptx"), (".doc", "docx")]:
        for path in Path(docs_dir).rglob(f"*{ext}"):
            target = path.with_suffix(f".{target_ext}")
            if target.exists():
                continue  # skip if already done
            print(f"Converting: {path.name} -> {target.name}")
            try:
                subprocess.run(
                    ["libreoffice", "--headless", "--convert-to", target_ext,
                     "--outdir", str(path.parent), str(path)],
                    check=True,
                    capture_output=True,
                )
                converted += 1
            except FileNotFoundError:
                print("  [ERROR] LibreOffice not found. Install with:")
                print("          sudo apt install libreoffice -y")
                raise SystemExit(1)
            except subprocess.CalledProcessError as e:
                print(f"  [WARN] Conversion failed for {path.name}: {e}")
    if converted:
        print(f"  Converted {converted} legacy file(s).\n")

    
def is_valid_ocr(text):
    if len(text.split()) < 5:
        return False

    return True


def extract_table_from_image(image_bytes):
    try:
        import os
        os.environ["PATH"] += r";C:\Program Files\Tesseract-OCR"
        ocr = TesseractOCR()
        doc = Img2TableImage(src=image_bytes)
        tables = doc.extract_tables(ocr=ocr, implicit_rows=True, borderless_tables=True)
        if not tables:
            return "", False
        rows = []
        for table in tables:
            df = table.df
            for _, row in df.iterrows():
                rows.append(" | ".join(str(c) for c in row if str(c).strip()))
        result = "\n".join(rows)
        return (result, True) if len(rows) >= 2 else ("", False)
    except Exception as e:
        print(f"[img2table ERROR]: {e}")
        return "", False


def extract_chart_data_deplot(image_bytes):
    try:
        image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
        inputs = processor(
            images=image,
            text="Generate underlying data table of the figure below:",
            return_tensors="pt"
        )
        predictions = model.generate(**inputs, max_new_tokens=512)
        result = processor.decode(predictions[0], skip_special_tokens=True)
        return result.strip() if result.strip() else None
    except Exception as e:
        print(f"[Deplot ERROR]: {e}")
        return None


def extract_chart_insights_from_image(image_bytes):
    """Use LLaVA to directly detect charts/graphs and extract insights — no OCR dependency."""
    try:
        image_b64 = base64.b64encode(image_bytes).decode("utf-8")
        prompt = """Analyze this image carefully.

If this image contains a chart, graph, or data visualization (bar chart, line graph, pie chart, scatter plot, etc.):
- Identify the chart type.
- Describe the axes, labels, and legend if visible.
- Summarize the key trend, comparison, or insight in 2-3 concise sentences.
- Note any significant data points or outliers.
- Return your analysis starting with: CHART_INSIGHT:

If this image does NOT contain a chart or graph, reply with exactly: NO_CHART"""

        response = requests.post(
            "http://localhost:11434/api/generate",
            json={
                "model": MULTIMODAL_MODEL,
                "prompt": prompt,
                "images": [image_b64],
                "stream": False,
                "options": {"temperature": 0, "num_predict": 300}
            },
            timeout=120
        )
        result = response.json().get("response", "").strip()

        if not result or "NO_CHART" in result.upper():
            return None

        if "CHART_INSIGHT:" in result:
            return result.split("CHART_INSIGHT:", 1)[1].strip()

        return result  # fallback: return whatever LLaVA gave

    except Exception as e:
        print(f"[LLaVA CHART ERROR]: {e}")
        return None


def detect_chart_image(text):
    if not text:
        return False

    lower = text.lower()
    keywords = [
        "chart", "graph", "x-axis", "y-axis", "legend", "series", "plot",
        "trend", "increase", "decrease", "percentage", "vs ", "axis",
    ]
    score = sum(1 for kw in keywords if kw in lower)
    return score >= 2 and bool(re.search(r"\d+", lower))


def summarize_chart_image(image_bytes, extracted_text=""):
    if not image_bytes:
        return None

    image_b64 = base64.b64encode(image_bytes).decode("utf-8")
    prompt = f"""
You are a data analyst. The user has provided a chart or graph image and the OCR text extracted from it.
Use the image and the text to summarize the main insight, trend, or comparison in 2-3 concise sentences.
If the chart information is unclear, explain what can be inferred and what is missing.

OCR text:
{extracted_text.strip()}

Instructions:
- Focus on the key trend or takeaway.
- Mention the axes, series, or direction if visible.
- Do not hallucinate details not present in the image or text.
"""
    try:
        response = requests.post(
            "http://localhost:11434/api/generate",
            json={
                "model": MULTIMODAL_MODEL,
                "prompt": prompt,
                "image": image_b64,
                "stream": False,
                "options": {"temperature": 0, "max_output_tokens": 200}
            },
            timeout=120
        )
        return response.json().get("response", "").strip()
    except Exception:
        return None


def detect_code(text):

    lines = text.split("\n")
    score = 0

    for line in lines:
        line = line.strip()

        # STRONG indicators only
        if (
            line.startswith("class ") or
            line.startswith("public ") or
            line.startswith("private ") or
            line.startswith("def ") or
            line.startswith("#include") or
            "::" in line or
            "->" in line or
            ("(" in line and ")" in line and "{" in line) or
            ("=" in line and ";" in line)
        ):
            score += 2   # strong signal

        # WEAK indicators
        elif (
            "{" in line or "}" in line or ";" in line
        ):
            score += 1

    return score >= 5

def detect_example(text):

    text = text.lower()

    patterns = [
        "example",
        "consider",
        "suppose",
        "let us",
        "case",
        "illustration"
    ]

    return any(p in text for p in patterns)

def classify_chunk_llm(text):

    prompt = f"""
You are a classifier.

Analyze the given text and classify it.

Return ONLY JSON:

{{
  "contains_code": true/false,
  "contains_example": true/false
}}

Rules:
- Code = programming structure, syntax, functions, class, etc.
- Example = explanation using a case, illustration, or scenario

Text:
{text[:800]}
"""

    response = requests.post(
        "http://localhost:11434/api/generate",
        json={
            "model": "llama3",
            "prompt": prompt,
            "stream": False,
            "options": {"temperature": 0}
        }
    )

    try:
        output = response.json()["response"]
        start = output.find("{")
        end = output.rfind("}") + 1
        result = json.loads(output[start:end])
        return result
    except:
        return {"contains_code": False, "contains_example": False}


# Text extraction — per format

def _extract_pptx(path: Path) -> list[dict]:
    """Extract per-slide content from a .pptx file."""
    prs = Presentation(str(path))
    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        # Title
        title = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text_frame.text.strip()

        # Body
        lines = []
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue

            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text_frame.text.strip()
                             for cell in row.cells
                             if cell.text_frame.text.strip()]
                    if cells:
                        lines.append(" | ".join(cells))
                lines.append("")   # blank line after table
                continue

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        lines.append("  " * para.level + t)
        body = "\n".join(lines)

        # Speaker notes
        notes = ""
        try:
            tf = slide.notes_slide.notes_text_frame
            notes = tf.text.strip() if tf else ""
        except Exception:
            pass

        parts = []
        if title:
            parts.append(f"Title: {title}")
        if body:
            parts.append(body)
        if notes:
            parts.append(f"Speaker Notes: {notes}")

        images = []

        # extract images from slide
        for shape in slide.shapes:
            if shape.shape_type == 13:  # picture
                try:
                    image_bytes = shape.image.blob
                    images.append(image_bytes)
                except Exception:
                    pass

        slides.append({
            "page_number": idx,
            "section":     title,
            "text":        normalize_text("\n\n".join(parts).strip()),
            "images":      images   
        })
    return slides


def _extract_pdf(path: Path) -> list[dict]:
    """Extract per-page content from a PDF using pdfplumber."""
    pages = []
    with pdfplumber.open(str(path)) as pdf:
        for idx, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""

            # Also extract tables as pipe-separated text
            tables = page.extract_tables()
            table_text = ""
            for table in tables:
                for row in table:
                    cleaned = [cell.strip() if cell else "" for cell in row]
                    table_text += " | ".join(cleaned) + "\n"

            combined = "\n\n".join(filter(None, [text.strip(), table_text.strip()]))
            pages.append({
                "page_number": idx,
                "section":     f"Page {idx}",
                "text":        normalize_text(combined.strip())
            })
    return pages

def _extract_docx_tables(doc) -> str:
    """Extract all tables from a docx as readable text."""
    table_texts = []
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            # Remove duplicate cells (merged cells repeat in python-docx)
            seen, unique = set(), []
            for c in cells:
                if c not in seen:
                    seen.add(c)
                    unique.append(c)
            rows.append(" | ".join(unique))
        table_texts.append("\n".join(rows))
    return "\n\n[TABLE]\n" + "\n\n[TABLE]\n".join(table_texts) if table_texts else ""

def _extract_docx(path: Path) -> list[dict]:
    """
    Extract content from a .docx file.
    Splits on Heading styles into logical sections.
    Falls back to fixed-size paragraph groups if no headings found.
    """
    doc      = DocxDocument(str(path))
    sections = []
    current_heading = "Introduction"
    current_paras   = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        style = para.style.name if para.style else ""
        is_heading = style.startswith("Heading") or style.startswith("Title")

        if is_heading:
            if current_paras:
                sections.append({
                    "page_number": len(sections) + 1,
                    "section":     current_heading,
                    "text":        normalize_text("\n".join(current_paras))
                })
            current_heading = text
            current_paras   = []
        else:
            current_paras.append(text)

    # Flush last section
    if current_paras:
        table_text = _extract_docx_tables(doc) if not sections else ""
        sections.append({
            "page_number": len(sections) + 1,
            "section":     current_heading,
            "text":       normalize_text("\n".join(current_paras)) + table_text,
        })

    # if no headings, treat as one big section
    if not sections:
        full_text = "\n".join(
            p.text.strip() for p in doc.paragraphs if p.text.strip()
        )
        full_text += _extract_docx_tables(doc)
        sections = [{"page_number": 1, "section": path.stem, "text": full_text}]

    return sections


def _extract_txt(path: Path) -> list[dict]:
    """
    Extract content from a plain .txt file.
    Splits on blank lines into paragraphs, groups into logical sections.
    """
    raw = path.read_text(encoding="utf-8", errors="ignore")

    # Split on double newlines (paragraph boundaries)
    paragraphs = [p.strip() for p in re.split(r"\n{2,}", raw) if p.strip()]

    # Group paragraphs into chunks of ~5 per section
    GROUP_SIZE = 5
    sections   = []
    for i in range(0, len(paragraphs), GROUP_SIZE):
        group = paragraphs[i : i + GROUP_SIZE]
        sections.append({
            "page_number": (i // GROUP_SIZE) + 1,
            "section":     f"Section {(i // GROUP_SIZE) + 1}",
            "text":        normalize_text("\n\n".join(group))
        })

    if not sections:
        sections = [{"page_number": 1, "section": path.stem, "text": raw.strip()}]

    return sections


def extract_content(path: Path) -> list[dict]:
    """Dispatch to the correct extractor based on file extension."""
    ext = path.suffix.lower()
    try:
        if ext == ".pptx":
            return _extract_pptx(path)
        elif ext == ".pdf":
            return _extract_pdf(path)
        elif ext == ".docx":
            return _extract_docx(path)
        elif ext == ".txt":
            return _extract_txt(path)
        else:
            print(f"  [WARN] Unsupported format skipped: {path.name}")
            return []
    except Exception as e:
        print(f"  [WARN] Failed to extract {path.name}: {e}")
        return []


# Chunking

def _split_text(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Split text into overlapping chunks by semantic, sentence-aware chunking."""

    if not text or len(text) <= size:
        return [text.strip()] if text.strip() else []

    sentences = sent_tokenize(text)
    chunks, current  = [], ""

    for sent in sentences:
        if len(current) + len(sent) + 1 <= size:
            current = (current + " " + sent).strip()
        else:
            if current:
                chunks.append(current)
            overlap_text = current[-overlap:] if len(current) > overlap else current
            current = (overlap_text + " " + sent).strip()

    if current:
        chunks.append(current)

    return chunks

def _split_text_semantic(text: str, threshold: float = SEMANTIC_SIM_THRESHOLD, size: int = CHUNK_SIZE) -> list[str]:
    """Semantic chunking — splits text when the topic shifts."""
    if not text or len(text) <= size:
        return [text.strip()] if text.strip() else []

    sentences = sent_tokenize(text)
    if len(sentences) <= 1:
        return [text.strip()]

    model      = SentenceTransformer(EMBED_MODEL)
    embeddings = model.encode(sentences, batch_size=64, show_progress_bar=False)

    chunks        = []
    current_sents = [sentences[0]]

    for i in range(1, len(sentences)):
        sim = cosine_similarity(
            embeddings[i - 1].reshape(1, -1),
            embeddings[i].reshape(1, -1)
        )[0][0]

        if sim < threshold:
            chunks.append(" ".join(current_sents))
            current_sents = [sentences[i]]
        else:
            current_sents.append(sentences[i])

    if current_sents:
        chunks.append(" ".join(current_sents))

    # Hard-split oversized chunks using sentence-aware fallback
    final_chunks = []
    for chunk in chunks:
        if len(chunk) > size:
            final_chunks.extend(_split_text(chunk, size=size))
        else:
            final_chunks.append(chunk)

    return [c for c in final_chunks if c.strip()]

def _split_text_old(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Split text into overlapping chunks, preferring sentence boundaries."""
    
    if not text or len(text) <= size:
        return [text] if text.strip() else []

    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks, current = [], ""

    for sent in sentences:
        if len(current) + len(sent) + 1 <= size:
            current = (current + " " + sent).strip()
        else:
            if current:
                chunks.append(current)
            while len(sent) > size:
                chunks.append(sent[:size])
                sent = sent[size - overlap:]
            current = sent

    if current:
        chunks.append(current)

    return chunks


def section_to_chunks(section: dict, file_name: str, file_path: str, file_type: str) -> list[Chunk]:
    """Convert one extracted section/page/slide into Chunk objects."""

    if CHUNKING_STRATEGY == "semantic":
        raw_chunks = _split_text_semantic(section["text"])
    elif CHUNKING_STRATEGY == "sentence":
        raw_chunks = _split_text(section["text"])
    else:
        raw_chunks = _split_text(section["text"])

    total = len(raw_chunks)
    result = []

    parent_str = f"{file_path}::p{section['page_number']}"
    parent_id  = hashlib.md5(parent_str.encode()).hexdigest()

    for idx, text in enumerate(raw_chunks):
        # Use full text hash + counter to guarantee uniqueness
        unique_str = f"{file_path}::p{section['page_number']}::c{idx}::{len(text)}::{text}"
        uid = hashlib.md5(unique_str.encode()).hexdigest()

        if detect_code(text):
            ctype = "code"
        elif re.search(r'[\$\\]|\\frac|\\sum|\\int|∫|∑|√', text):
            ctype = "equation"
        else:
            ctype = "text"

        result.append(Chunk(
            doc_id      = uid,
            text        = text,
            file_name   = file_name,
            file_path   = file_path,
            file_type   = file_type,
            page_number = section["page_number"],
            section     = section["section"],
            is_image    = False,
            image_type  = "text",
            chunk_index = idx,
            total_chunks= total,
            parent_id   = parent_id,
            chunk_type  = ctype
        ))
    return result


# Vector DB (ChromaDB)

def store_in_vector_db(chunks: list[Chunk]) -> chromadb.Collection:
    """Embed all chunks and upsert into ChromaDB."""

    # Deduplicate by doc_id before doing anything
    seen_ids = set()
    unique_chunks = []
    for c in chunks:
        if c.doc_id not in seen_ids:
            seen_ids.add(c.doc_id)
            unique_chunks.append(c)
        else:
            print(f"  [DEDUP] Skipping duplicate id: {c.doc_id}")
    chunks = unique_chunks

    print(f"\n[VectorDB] Loading embedding model: {EMBED_MODEL}")
    model = SentenceTransformer(EMBED_MODEL)

    print(f"[VectorDB] Embedding {len(chunks)} chunks ...")
    texts      = [c.text for c in chunks]
    embeddings = model.encode(texts, show_progress_bar=True, batch_size=64)

    client = chromadb.PersistentClient(
        path=CHROMA_DB_DIR,
        settings=Settings(anonymized_telemetry=False),
    )
    collection = client.get_or_create_collection(
        name=COLLECTION_NAME,
        metadata={"hnsw:space": "cosine"},
    )

    BATCH = 500
    print(f"[VectorDB] Upserting into '{COLLECTION_NAME}' ...")
    for i in tqdm(range(0, len(chunks), BATCH), desc="Batches"):

        batch = chunks[i : i + BATCH]

        metadatas = []

        for c in batch:

            contains_code = (c.chunk_type == "code")
            contains_example = detect_example(c.text)

            if  len(c.text) > 60:

                classification = classify_chunk_llm(c.text)

                # safer access (avoid crashes)
                contains_code = contains_code or classification.get("contains_code", False)
                contains_example = contains_example or classification.get("contains_example", False)

            # ---------- STEP 3: BUILD METADATA ----------
            meta = {
                "file_name":    c.file_name,
                "file_path":    c.file_path,
                "file_type":    c.file_type,
                "page_number":  c.page_number,
                "section":      c.section,
                "chunk_index":  c.chunk_index,
                "total_chunks": c.total_chunks,

                "source": "image" if c.is_image == True else "text",
                "image_type":   c.image_type,

                "contains_code": contains_code,
                "contains_example": contains_example,
                "parent_id":    c.parent_id,      
                "chunk_type":   c.chunk_type
            }

            metadatas.append(meta)

        # ---------- UPSERT ----------
        collection.upsert(
            ids        = [c.doc_id for c in batch],
            documents  = [c.text for c in batch],
            embeddings = embeddings[i : i + BATCH].tolist(),
            metadatas  = metadatas
        )

    print(f"[VectorDB] {collection.count()} total documents in collection.\n")
    return collection


# Main pipeline

def reconstruct_code_llm(text):

    prompt = f"""
The following is OCR-extracted code and may be corrupted.

Fix syntax, formatting, and structure.

Return ONLY corrected code.
Dont give the changes. Return only the code.

{text}
"""

    response = requests.post(
        "http://localhost:11434/api/generate",
        json={
            "model": "llama3",
            "prompt": prompt,
            "stream": False,
            "options": {"temperature": 0}
        }
    )

    try:
        return response.json()["response"].strip()
    except:
        return text

def run_pipeline(docs_dir: str = DOCS_DIR) -> chromadb.Collection:
    print(f"\n{'='*58}")
    print(f"  Generalized Knowledge Base Construction")
    print(f"{'='*58}")
    print(f"  Source dir : {docs_dir}")
    print(f"  Formats    : {', '.join(sorted(SUPPORTED_EXTS))}")
    print(f"  ChromaDB   : {CHROMA_DB_DIR}")
    print(f"  Collection : {COLLECTION_NAME}")
    print(f"{'='*58}\n")

    # Step 1 — convert legacy formats
    print("[Step 1] Converting legacy .ppt / .doc files ...")
    convert_legacy_files(docs_dir)

    # Step 2 — discover all supported files
    all_files = sorted(
        p for p in Path(docs_dir).rglob("*")
        if p.suffix.lower() in SUPPORTED_EXTS and p.is_file()
    )

    if not all_files:
        raise FileNotFoundError(f"No supported files found in: {docs_dir}")

    # Summary by type
    from collections import Counter
    type_counts = Counter(p.suffix.lower() for p in all_files)
    print(f"[Step 2] Found {len(all_files)} file(s):")
    for ext, count in sorted(type_counts.items()):
        print(f"         {ext:6s}  ->  {count} file(s)")
    print()

    # Step 3 — extract + chunk
    all_chunks: list[Chunk] = []
    print("[Step 3] Extracting and chunking ...")
    for path in tqdm(all_files, desc="Files"):
        file_type = path.suffix.lower().lstrip(".")
        sections  = extract_content(path)
        file_chunks = []
        for section in sections:
            if not section["text"].strip():
                continue
            # ---------- TEXT CHUNKS ----------
            file_chunks.extend(
                section_to_chunks(section, path.name, str(path.resolve()), file_type)
            )

            # ---------- IMAGE CHUNKS ----------
            images = section.get("images", [])

            for img_bytes in images:
                image_type = classify_image_type(img_bytes)

                processed_text = ""

                # Reconstruct code using LLM only if it's a code image
                if image_type == "code":
                    code_text = extract_code_from_image(img_bytes)
                    processed_text = reconstruct_code_llm(code_text)

                # --- Step 1: Try table extraction via Img2Table ---
                if not processed_text:
                    table_text, is_table = extract_table_from_image(img_bytes)
                    if is_table and table_text:
                        processed_text = table_text
                        image_type = "table"

                # --- Step 2: Try chart/graph insight extraction via deplot ---
                if not processed_text:
                    # chart_insight = extract_chart_insights_from_image(img_bytes)
                    chart_insight = extract_chart_data_deplot(img_bytes)
                    if chart_insight:
                        processed_text = f"Chart insights:\n{chart_insight}"
                        image_type = "chart"

                # --- Step 3: Fallback to OCR for plain text images ---
                if not processed_text:
                    processed_text = ocr_image(img_bytes)
                    image_type = "ocr"

                if not is_valid_ocr(processed_text):
                    continue

                uid = hashlib.md5(
                    f"{path}::img::{section['page_number']}::{processed_text}".encode()
                ).hexdigest()

                print(f"    {image_type, processed_text}")

                file_chunks.append(Chunk(
                    doc_id       = uid,
                    text         = processed_text,
                    file_name    = path.name,
                    file_path    = str(path.resolve()),
                    file_type    = file_type,
                    page_number  = section["page_number"],
                    section      = section["section"],
                    is_image     = True,
                    image_type   = image_type,
                    chunk_index  = -1,
                    total_chunks = -1,
                    parent_id    = hashlib.md5(f"{path}::p{section['page_number']}".encode()).hexdigest(),
                    # chunk_type   = img_chunk_type,   # "table", "chart", "code", "image"
                ))
        tqdm.write(
            f"  {path.name:45s}  {len(sections):4d} sections  ->  {len(file_chunks):4d} chunks"
        )
        all_chunks.extend(file_chunks)

    if not all_chunks:
        raise ValueError("No text could be extracted from any file.")

    print(f"\n[Step 3] Total chunks: {len(all_chunks)}\n")

    # Step 4 — embed + store
    print("[Step 4] Vectorizing and storing in ChromaDB ...")
    collection = store_in_vector_db(all_chunks)

    print(" Knowledge base is ready.\n")
    return collection


if __name__ == "__main__":

    collection = run_pipeline()

    print("\n========== DEBUG: VERIFYING DATA ==========\n")

    # Fetch stored data
    data = collection.get(include=["documents", "metadatas"])

    docs = data["documents"]
    metas = data["metadatas"]

    print(f"Total stored chunks: {len(docs)}\n")

    # Count types
    image_indices = [i for i, m in enumerate(metas) if m.get("source") == "image"]
    text_indices  = [i for i, m in enumerate(metas) if m.get("source") == "text"]
    code_chunks = [i for i, m in enumerate(metas) if m.get("chunk_type") == "code"]
    table_chunks = [i for i, m in enumerate(metas) if m.get("chunk_type") == "table"]
    chart_chunks = [i for i, m in enumerate(metas) if m.get("chunk_type") == "chart"]
    other_chunks = [i for i, m in enumerate(metas) if m.get("chunk_type") == "image"]

    print(f"Text chunks  : {len(text_indices)}")
    print(f"Image chunks : {len(image_indices)}\n")
    print(f"Code chunks : {len(code_chunks)}\n")
    print(f"Table chunks : {len(table_chunks)}\n")
    print(f"Chart chunks : {len(chart_chunks)}\n")
    print(f"Other chunks : {len(other_chunks)}\n")

    # Code / example stats
    code_count = sum(1 for m in metas if m.get("contains_code"))
    example_count = sum(1 for m in metas if m.get("contains_example"))

    print(f"Code chunks     : {code_count}")
    print(f"Example chunks  : {example_count}\n")

    # ---------- VIEW IMAGE CHUNKS ----------
    print("\n========== IMAGE CHUNKS (OCR OUTPUT) ==========\n")

    if not image_indices:
        print("No image chunks found ❌\n")
    else:
        for idx in image_indices[:5]:   # show first 5 images
            print(f"\n--- Image Chunk {idx} ---")
            print("TEXT:\n", docs[idx][:500])   # OCR content
            print("META:", metas[idx])

    # ---------- VIEW CODE CHUNKS ----------
    print("\n========== CODE CHUNKS ==========\n")

    code_indices = [i for i, m in enumerate(metas) if m.get("contains_code")]

    for idx in code_indices[:5]:
        print(f"\n--- Code Chunk {idx} ---")
        print("TEXT:\n", docs[idx][:500])
        print("META:", metas[idx])

    # ---------- VIEW RANDOM SAMPLE ----------
    print("\n========== RANDOM SAMPLE ==========\n")

    for i in range(min(5, len(docs))):
        print(f"\n--- Chunk {i} ---")
        print("TEXT:\n", docs[i][:300])
        print("META:", metas[i])